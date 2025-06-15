import os
import shutil
from typing import List, Dict, Optional
from pathlib import Path
from pptx import Presentation
from src.models import ProcessedSlide, Template
from src.template_manager import TemplateManager

class PowerPointGenerator:
    """PowerPoint生成クラス"""

    def __init__(self, template_manager: TemplateManager):
        self.template_manager = template_manager
        self.output_dir = Path("output")
        self.output_dir.mkdir(exist_ok=True)
        self.tmp_dir = Path("tmp_pptx")
        self.tmp_dir.mkdir(exist_ok=True)

    def generate_presentation(self, processed_slides: List[ProcessedSlide], output_filename: str = "generated_slides.pptx") -> str:
        """処理済みスライドからPowerPointプレゼンテーションを生成"""

        # 新しいプレゼンテーションを作成
        prs = Presentation()

        # 最初のデフォルトスライドを削除
        if len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

        for processed_slide in processed_slides:
            self._add_slide_from_template(prs, processed_slide)

        # 出力ファイルパス
        output_path = self.output_dir / output_filename

        # プレゼンテーションを保存
        prs.save(str(output_path))

        return str(output_path)

    def _add_slide_from_template(self, prs: Presentation, processed_slide: ProcessedSlide) -> None:
        """テンプレートからスライドを追加（テンプレートpptxを一時ファイルに複製してから編集）"""

        # テンプレートのPowerPointファイルパスを取得
        template_pptx_path = self.template_manager.get_template_pptx_path(processed_slide.template_name)

        if not template_pptx_path or not template_pptx_path.exists():
            print(f"Warning: Template file not found for {processed_slide.template_name}")
            return

        # 一時ファイル名を作成
        tmp_pptx_path = self.tmp_dir / f"tmp_{processed_slide.template_name}_{processed_slide.slide_number}.pptx"
        shutil.copy(str(template_pptx_path), str(tmp_pptx_path))

        # 一時ファイルを開く
        template_prs = Presentation(str(tmp_pptx_path))

        if len(template_prs.slides) == 0:
            print(f"Warning: No slides found in template {processed_slide.template_name}")
            return

        # テンプレートの最初のスライドを取得
        template_slide = template_prs.slides[0]

        # 各shapeのテキストを置換
        for shape in template_slide.shapes:
            if hasattr(shape, 'text'):
                original_text = shape.text.strip()
                if original_text in processed_slide.content_mapping:
                    shape.text = processed_slide.content_mapping[original_text]

        # 編集したスライドをメインプレゼンテーションに追加
        self._append_slide_from_other_presentation(prs, template_slide)

        # 一時ファイルを削除（必要に応じて）
        try:
            tmp_pptx_path.unlink()
        except Exception:
            pass

    def _append_slide_from_other_presentation(self, prs: Presentation, slide_to_copy) -> None:
        """
        別のPresentationのスライドを現在のPresentationに追加する。
        pptxの仕様上、直接スライドをコピーするAPIはないため、shapesを新規スライドに複製する。
        """
        blank_layout = prs.slide_layouts[6]  # Blank layout
        new_slide = prs.slides.add_slide(blank_layout)
        for shape in slide_to_copy.shapes:
            self._copy_shape(new_slide, shape)

    def _copy_shape(self, slide, template_shape) -> None:
        """shapeを新しいスライドにコピー（テキストボックスのみ対応、画像等はプレースホルダー）"""
        left = template_shape.left
        top = template_shape.top
        width = template_shape.width
        height = template_shape.height

        if hasattr(template_shape, 'text'):
            new_textbox = slide.shapes.add_textbox(left, top, width, height)
            new_textbox.text = template_shape.text
        else:
            placeholder_textbox = slide.shapes.add_textbox(left, top, width, height)
            placeholder_textbox.text = "[要素]"

    def generate_single_slide(self, processed_slide: ProcessedSlide, output_filename: str = None) -> str:
        """単一スライドのPowerPointファイルを生成（テンプレートpptxを一時ファイルに複製して編集）"""

        if output_filename is None:
            output_filename = f"slide_{processed_slide.slide_number}_{processed_slide.template_name}.pptx"

        # テンプレートのPowerPointファイルパスを取得
        template_pptx_path = self.template_manager.get_template_pptx_path(processed_slide.template_name)

        if not template_pptx_path or not template_pptx_path.exists():
            raise FileNotFoundError(f"Template file not found: {processed_slide.template_name}")

        # 一時ファイル名を作成
        tmp_pptx_path = self.tmp_dir / f"tmp_{processed_slide.template_name}_{processed_slide.slide_number}.pptx"
        shutil.copy(str(template_pptx_path), str(tmp_pptx_path))

        # 一時ファイルを開く
        template_prs = Presentation(str(tmp_pptx_path))

        if len(template_prs.slides) == 0:
            raise ValueError(f"No slides found in template: {processed_slide.template_name}")

        # テンプレートの最初のスライドを取得
        template_slide = template_prs.slides[0]

        # 各shapeのテキストを置換
        for shape in template_slide.shapes:
            if hasattr(shape, 'text'):
                original_text = shape.text.strip()
                if original_text in processed_slide.content_mapping:
                    shape.text = processed_slide.content_mapping[original_text]

        # 出力ファイルパス
        output_path = self.output_dir / output_filename

        # プレゼンテーションを保存
        template_prs.save(str(output_path))

        # 一時ファイルを削除（必要に応じて）
        try:
            tmp_pptx_path.unlink()
        except Exception:
            pass

        return str(output_path)
