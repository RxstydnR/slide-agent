"""
PowerPointテンプレートファイル作成用ヘルパースクリプト

各テンプレートフォルダに基本的なPowerPointテンプレートファイルを作成します。
実際の運用では、これらのファイルをPowerPointで手動編集してデザインを調整してください。
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide_template():
    """タイトルスライドテンプレートを作成"""
    prs = Presentation()
    
    # 最初のスライドを削除
    if len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    # 空のレイアウトでスライドを追加
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # タイトル用テキストボックス
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_box.text = "title"
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    
    # サブタイトル用テキストボックス
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
    subtitle_box.text = "subtitle"
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(24)
    
    # 発表者用テキストボックス
    presenter_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.6))
    presenter_box.text = "presenter"
    presenter_frame = presenter_box.text_frame
    presenter_para = presenter_frame.paragraphs[0]
    presenter_para.alignment = PP_ALIGN.CENTER
    presenter_para.font.size = Pt(18)
    
    # 日付用テキストボックス
    date_box = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(8), Inches(0.6))
    date_box.text = "date"
    date_frame = date_box.text_frame
    date_para = date_frame.paragraphs[0]
    date_para.alignment = PP_ALIGN.CENTER
    date_para.font.size = Pt(16)
    
    return prs

def create_single_column_template():
    """1カラムテキストテンプレートを作成"""
    prs = Presentation()
    
    # 最初のスライドを削除
    if len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    # 空のレイアウトでスライドを追加
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # タイトル用テキストボックス
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_box.text = "title"
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    
    # 本文用テキストボックス
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(5))
    content_box.text = "main_text"
    content_frame = content_box.text_frame
    content_para = content_frame.paragraphs[0]
    content_para.font.size = Pt(18)
    
    return prs

def create_lead_text_template():
    """リード文+1カラムテキストテンプレートを作成"""
    prs = Presentation()
    
    # 最初のスライドを削除
    if len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    # 空のレイアウトでスライドを追加
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # タイトル用テキストボックス
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_box.text = "title"
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    
    # リード文用テキストボックス
    lead_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1.2))
    lead_box.text = "lead"
    lead_frame = lead_box.text_frame
    lead_para = lead_frame.paragraphs[0]
    lead_para.font.size = Pt(20)
    lead_para.font.bold = True
    
    # 本文用テキストボックス
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(3.5))
    content_box.text = "main_text"
    content_frame = content_box.text_frame
    content_para = content_frame.paragraphs[0]
    content_para.font.size = Pt(16)
    
    return prs

def create_two_column_template():
    """2カラム（画像＋テキスト）テンプレートを作成"""
    prs = Presentation()
    
    # 最初のスライドを削除
    if len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    # 空のレイアウトでスライドを追加
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # タイトル用テキストボックス
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_box.text = "title"
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    
    # リード文用テキストボックス
    lead_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(0.8))
    lead_box.text = "lead"
    lead_frame = lead_box.text_frame
    lead_para = lead_frame.paragraphs[0]
    lead_para.font.size = Pt(18)
    lead_para.font.bold = True
    
    # 左側（画像エリア）用テキストボックス
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(4.2), Inches(4))
    left_box.text = "left_image"
    left_frame = left_box.text_frame
    left_para = left_frame.paragraphs[0]
    left_para.alignment = PP_ALIGN.CENTER
    left_para.font.size = Pt(16)
    
    # 右側（テキストエリア）用テキストボックス
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(2.8), Inches(4.2), Inches(4))
    right_box.text = "right_text"
    right_frame = right_box.text_frame
    right_para = right_frame.paragraphs[0]
    right_para.font.size = Pt(16)
    
    return prs

def create_end_slide_template():
    """エンドスライドテンプレートを作成"""
    prs = Presentation()
    
    # 最初のスライドを削除
    if len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    # 空のレイアウトでスライドを追加
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # エンドメッセージ用テキストボックス
    end_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    end_box.text = "end_message"
    end_frame = end_box.text_frame
    end_para = end_frame.paragraphs[0]
    end_para.alignment = PP_ALIGN.CENTER
    end_para.font.size = Pt(36)
    end_para.font.bold = True
    
    return prs

def main():
    """メイン関数"""
    templates_dir = Path("templates")
    
    # テンプレート作成関数のマッピング
    template_creators = {
        "タイトルスライド": create_title_slide_template,
        "1カラムテキスト": create_single_column_template,
        "リード文+1カラムテキスト": create_lead_text_template,
        "2カラム（画像＋テキスト）": create_two_column_template,
        "エンドスライド": create_end_slide_template,
    }
    
    for template_name, creator_func in template_creators.items():
        template_dir = templates_dir / template_name
        if template_dir.exists():
            pptx_path = template_dir / "template.pptx"
            
            print(f"Creating template: {template_name}")
            prs = creator_func()
            prs.save(str(pptx_path))
            print(f"Saved: {pptx_path}")
        else:
            print(f"Warning: Template directory not found: {template_dir}")
    
    print("\nテンプレートファイルの作成が完了しました。")
    print("実際の運用では、PowerPointでこれらのファイルを開いて")
    print("デザインやレイアウトを調整してください。")

if __name__ == "__main__":
    main()
